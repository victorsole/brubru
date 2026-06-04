"""
User Documents API Router

FastAPI endpoints for managing user documents in My EU Bubble.
Part of My EU Bubble - Phase 2: User Content Repository

Endpoints:
- POST /api/documents - Create document
- GET /api/documents - List user's documents
- GET /api/documents/{id} - Get specific document
- PATCH /api/documents/{id} - Update document
- DELETE /api/documents/{id} - Delete document
- POST /api/documents/{id}/version - Create new version
- GET /api/documents/{id}/versions - Get document versions
- GET /api/documents/stats - Get document statistics
"""

import logging
import re
from typing import List, Optional
from datetime import datetime
from uuid import UUID

from fastapi import APIRouter, HTTPException, Query, Depends, status
from fastapi.responses import StreamingResponse
from sqlalchemy.orm import Session
from sqlalchemy import select, and_, or_, func

from models.user import User
from models.user_document import UserDocument
from models.amendment import Amendment
from services.document_generation.document_exporter import export_docx, export_pdf
from schemas.user_document_schemas import (
    UserDocumentCreate,
    UserDocumentUpdate,
    UserDocumentResponse,
    UserDocumentListResponse,
    DocumentFilterParams,
    DocumentStatsResponse,
    DocumentVersionResponse,
    AmendmentCreate,
    AnalysisCreate,
    StrategyCreate,
    NoteCreate,
)
from core.database import get_db
from api.auth_optional import get_current_user_dev as get_current_user

logger = logging.getLogger(__name__)

router = APIRouter(
    prefix="/api/documents",
    tags=["User Documents"],
    responses={404: {"description": "Not found"}}
)


# ============================================================================
# Document CRUD Operations
# ============================================================================

@router.post(
    "",
    response_model=UserDocumentResponse,
    status_code=status.HTTP_201_CREATED,
    summary="Create document",
    description="Create a new document (amendment, analysis, strategy, or note)"
)
async def create_document(
    document: UserDocumentCreate,
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db)
) -> UserDocumentResponse:
    """
    Create a new user document.

    Automatically extracts policy areas from content if not provided.
    """
    try:
        # Create document
        new_doc = UserDocument(
            user_id=current_user.id,
            document_type=document.document_type,
            title=document.title,
            content=document.content,
            celex_number=document.celex_number,
            procedure_reference=document.procedure_reference,
            policy_areas=document.policy_areas,
            tags=document.tags,
            is_private=document.is_private,
            # Type-specific fields
            amendment_target=document.amendment_target if hasattr(document, 'amendment_target') else None,
            amendment_justification=document.amendment_justification if hasattr(document, 'amendment_justification') else None,
            amendment_status=document.amendment_status if hasattr(document, 'amendment_status') else None,
            executive_summary=document.executive_summary if hasattr(document, 'executive_summary') else None,
            methodology=document.methodology if hasattr(document, 'methodology') else None,
            conclusions=document.conclusions if hasattr(document, 'conclusions') else None,
            recommendations=document.recommendations if hasattr(document, 'recommendations') else None,
            objectives=document.objectives if hasattr(document, 'objectives') else None,
            action_plan=document.action_plan if hasattr(document, 'action_plan') else None,
            timeline=document.timeline if hasattr(document, 'timeline') else None,
            stakeholders=document.stakeholders if hasattr(document, 'stakeholders') else None,
            doc_metadata=document.doc_metadata if hasattr(document, 'doc_metadata') else None,
        )

        db.add(new_doc)
        db.commit()
        db.refresh(new_doc)

        logger.info(f"User {current_user.email} created document: {new_doc.title}")

        # Create response with computed fields
        response = UserDocumentResponse.from_orm(new_doc)
        response.word_count = new_doc.word_count

        return response

    except Exception as e:
        db.rollback()
        logger.error(f"Failed to create document: {str(e)}")
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=f"Failed to create document: {str(e)}"
        )


@router.get(
    "",
    response_model=UserDocumentListResponse,
    summary="List documents",
    description="Get user's documents with filtering and pagination"
)
async def list_documents(
    document_type: Optional[str] = Query(None, description="Filter by type"),
    policy_areas: Optional[List[str]] = Query(None, description="Filter by policy areas"),
    tags: Optional[List[str]] = Query(None, description="Filter by tags"),
    search: Optional[str] = Query(None, max_length=200, description="Search in title and content"),
    limit: int = Query(50, ge=1, le=200),
    offset: int = Query(0, ge=0),
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db)
) -> UserDocumentListResponse:
    """List user's documents with filters"""
    try:
        # Build query
        query = db.query(UserDocument).filter(UserDocument.user_id == current_user.id)

        # Apply filters
        if document_type:
            query = query.filter(UserDocument.document_type == document_type)

        if policy_areas:
            # Filter documents that have any of the specified policy areas
            for area in policy_areas:
                query = query.filter(UserDocument.policy_areas.contains([area]))

        if tags:
            # Filter documents that have any of the specified tags
            for tag in tags:
                query = query.filter(UserDocument.tags.contains([tag]))

        if search:
            # Search in title and content
            search_pattern = f"%{search}%"
            query = query.filter(
                or_(
                    UserDocument.title.ilike(search_pattern),
                    UserDocument.content.ilike(search_pattern)
                )
            )

        # Order by updated_at
        query = query.order_by(UserDocument.updated_at.desc())

        # Get total count
        total = query.count()

        # Apply pagination
        documents = query.offset(offset).limit(limit).all()

        # Build response with computed fields
        doc_responses = []
        for doc in documents:
            doc_response = UserDocumentResponse.from_orm(doc)
            doc_response.word_count = doc.word_count
            doc_responses.append(doc_response)

        return UserDocumentListResponse(
            documents=doc_responses,
            total=total,
            limit=limit,
            offset=offset,
            has_more=(offset + len(documents)) < total
        )

    except Exception as e:
        logger.error(f"Failed to list documents: {str(e)}")
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=f"Failed to list documents: {str(e)}"
        )


@router.get(
    "/{document_id}",
    response_model=UserDocumentResponse,
    summary="Get document",
    description="Get specific document by ID"
)
async def get_document(
    document_id: UUID,
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db)
) -> UserDocumentResponse:
    """Get specific document"""
    try:
        document = db.query(UserDocument).filter(
            and_(
                UserDocument.id == document_id,
                UserDocument.user_id == current_user.id
            )
        ).first()

        if not document:
            raise HTTPException(
                status_code=status.HTTP_404_NOT_FOUND,
                detail="Document not found"
            )

        # Increment view count
        document.increment_view_count()
        db.commit()

        # Build response
        response = UserDocumentResponse.from_orm(document)
        response.word_count = document.word_count

        return response

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Failed to get document: {str(e)}")
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=f"Failed to get document: {str(e)}"
        )


@router.patch(
    "/{document_id}",
    response_model=UserDocumentResponse,
    summary="Update document",
    description="Update document fields"
)
async def update_document(
    document_id: UUID,
    updates: UserDocumentUpdate,
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db)
) -> UserDocumentResponse:
    """Update document"""
    try:
        document = db.query(UserDocument).filter(
            and_(
                UserDocument.id == document_id,
                UserDocument.user_id == current_user.id
            )
        ).first()

        if not document:
            raise HTTPException(
                status_code=status.HTTP_404_NOT_FOUND,
                detail="Document not found"
            )

        # Update fields
        update_data = updates.dict(exclude_unset=True)
        for field, value in update_data.items():
            setattr(document, field, value)

        db.commit()
        db.refresh(document)

        logger.info(f"User {current_user.email} updated document: {document.title}")

        # Build response
        response = UserDocumentResponse.from_orm(document)
        response.word_count = document.word_count

        return response

    except HTTPException:
        raise
    except Exception as e:
        db.rollback()
        logger.error(f"Failed to update document: {str(e)}")
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=f"Failed to update document: {str(e)}"
        )


@router.delete(
    "/{document_id}",
    status_code=status.HTTP_204_NO_CONTENT,
    summary="Delete document",
    description="Delete document permanently"
)
async def delete_document(
    document_id: UUID,
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    """Delete document"""
    try:
        document = db.query(UserDocument).filter(
            and_(
                UserDocument.id == document_id,
                UserDocument.user_id == current_user.id
            )
        ).first()

        if not document:
            raise HTTPException(
                status_code=status.HTTP_404_NOT_FOUND,
                detail="Document not found"
            )

        # If uploaded document, also clean up filesystem storage
        if document.document_type == 'uploaded' and document.storage_document_id:
            try:
                from services.storage.document_storage import get_document_storage
                doc_storage = get_document_storage()
                doc_storage.delete_document(document.storage_document_id)
                logger.info(f"Cleaned up filesystem for upload {document.storage_document_id}")
            except Exception as storage_err:
                logger.warning(f"Failed to clean up filesystem storage: {storage_err}")

        db.delete(document)
        db.commit()

        logger.info(f"User {current_user.email} deleted document: {document.title}")

        return None

    except HTTPException:
        raise
    except Exception as e:
        db.rollback()
        logger.error(f"Failed to delete document: {str(e)}")
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=f"Failed to delete document: {str(e)}"
        )


# ============================================================================
# Version Control
# ============================================================================

@router.post(
    "/{document_id}/version",
    response_model=UserDocumentResponse,
    status_code=status.HTTP_201_CREATED,
    summary="Create new version",
    description="Create a new version of the document"
)
async def create_document_version(
    document_id: UUID,
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db)
) -> UserDocumentResponse:
    """Create new version of document"""
    try:
        document = db.query(UserDocument).filter(
            and_(
                UserDocument.id == document_id,
                UserDocument.user_id == current_user.id
            )
        ).first()

        if not document:
            raise HTTPException(
                status_code=status.HTTP_404_NOT_FOUND,
                detail="Document not found"
            )

        # Create new version
        new_version = document.create_new_version(db)
        db.commit()
        db.refresh(new_version)

        logger.info(f"User {current_user.email} created version {new_version.version} of document: {document.title}")

        # Build response
        response = UserDocumentResponse.from_orm(new_version)
        response.word_count = new_version.word_count

        return response

    except HTTPException:
        raise
    except Exception as e:
        db.rollback()
        logger.error(f"Failed to create document version: {str(e)}")
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=f"Failed to create version: {str(e)}"
        )


@router.get(
    "/{document_id}/versions",
    response_model=List[DocumentVersionResponse],
    summary="Get document versions",
    description="Get all versions of a document"
)
async def get_document_versions(
    document_id: UUID,
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db)
) -> List[DocumentVersionResponse]:
    """Get all versions of document"""
    try:
        # Get the document
        document = db.query(UserDocument).filter(
            and_(
                UserDocument.id == document_id,
                UserDocument.user_id == current_user.id
            )
        ).first()

        if not document:
            raise HTTPException(
                status_code=status.HTTP_404_NOT_FOUND,
                detail="Document not found"
            )

        # Get all versions (documents with this parent_document_id or this document itself)
        versions = db.query(UserDocument).filter(
            and_(
                UserDocument.user_id == current_user.id,
                or_(
                    UserDocument.id == document_id,
                    UserDocument.parent_document_id == document_id
                )
            )
        ).order_by(UserDocument.version.asc()).all()

        return [DocumentVersionResponse.from_orm(v) for v in versions]

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Failed to get document versions: {str(e)}")
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=f"Failed to get versions: {str(e)}"
        )


# ============================================================================
# Statistics
# ============================================================================

@router.get(
    "/stats/summary",
    response_model=DocumentStatsResponse,
    summary="Get document statistics",
    description="Get statistics about user's documents"
)
async def get_document_stats(
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db)
) -> DocumentStatsResponse:
    """Get document statistics"""
    try:
        # Total documents
        total = db.query(UserDocument).filter(
            UserDocument.user_id == current_user.id
        ).count()

        # By type
        by_type_query = db.query(
            UserDocument.document_type,
            func.count(UserDocument.id).label('count')
        ).filter(
            UserDocument.user_id == current_user.id
        ).group_by(UserDocument.document_type).all()

        by_type = {doc_type: count for doc_type, count in by_type_query}

        # By policy area
        policy_area_query = db.query(
            func.unnest(UserDocument.policy_areas).label('policy_area'),
            func.count().label('count')
        ).filter(
            UserDocument.user_id == current_user.id
        ).group_by('policy_area').order_by(func.count().desc()).limit(10).all()

        by_policy_area = {area: count for area, count in policy_area_query if area}

        # Total word count
        documents = db.query(UserDocument).filter(
            UserDocument.user_id == current_user.id
        ).all()
        total_words = sum(doc.word_count for doc in documents)

        # Most viewed
        most_viewed = db.query(UserDocument).filter(
            UserDocument.user_id == current_user.id
        ).order_by(UserDocument.view_count.desc()).limit(5).all()

        most_viewed_responses = []
        for doc in most_viewed:
            doc_response = UserDocumentResponse.from_orm(doc)
            doc_response.word_count = doc.word_count
            most_viewed_responses.append(doc_response)

        # Count amendments from the amendments table (not user_document)
        logger.info(f"📊 Getting stats for user: {current_user.email} (ID: {current_user.id})")
        total_amendments = db.query(Amendment).filter(
            Amendment.user_id == current_user.id
        ).count()
        logger.info(f"📊 Found {total_amendments} amendments for user {current_user.email}")

        return DocumentStatsResponse(
            total_documents=total,
            by_type=by_type,
            by_policy_area=by_policy_area,
            total_amendments=total_amendments,
            total_analyses=by_type.get('analysis', 0),
            total_strategies=by_type.get('strategy', 0),
            total_notes=by_type.get('note', 0),
            total_uploaded=by_type.get('uploaded', 0),
            total_word_count=total_words,
            most_viewed_documents=most_viewed_responses
        )

    except Exception as e:
        logger.error(f"Failed to get document stats: {str(e)}")
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=f"Failed to get statistics: {str(e)}"
        )


# ============================================================================
# Export (DOCX / PDF)
# ============================================================================

_SLUG_RE = re.compile(r"[^A-Za-z0-9._-]+")


def _safe_filename(title: str, ext: str) -> str:
    base = _SLUG_RE.sub("_", (title or "document").strip())[:80].strip("_") or "document"
    return f"{base}.{ext}"


@router.get(
    "/{document_id}/export",
    summary="Export document",
    description="Download a saved document as DOCX or PDF"
)
async def export_document(
    document_id: UUID,
    fmt: str = Query("docx", alias="format", pattern="^(docx|pdf)$"),
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db),
):
    """Stream a user document as a DOCX or PDF attachment."""
    document = db.query(UserDocument).filter(
        and_(
            UserDocument.id == document_id,
            UserDocument.user_id == current_user.id,
        )
    ).first()

    if not document:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Document not found")

    title = document.title or "Document"
    content = document.content or ""

    try:
        if fmt == "docx":
            payload = export_docx(title, content)
            media = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
            filename = _safe_filename(title, "docx")
        else:
            payload = export_pdf(title, content)
            media = "application/pdf"
            filename = _safe_filename(title, "pdf")
    except Exception as exc:
        logger.exception("Failed to render export for %s", document_id)
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=f"Failed to export document: {exc}",
        ) from exc

    import io

    return StreamingResponse(
        io.BytesIO(payload),
        media_type=media,
        headers={"Content-Disposition": f'attachment; filename="{filename}"'},
    )


# ---------------------------------------------------------------------------
# Presentation export: markdown slides -> PPTX
# ---------------------------------------------------------------------------

def _render_pptx_from_slides_markdown(title: str, markdown: str) -> bytes:
    """
    Parse the slide-by-slide markdown produced by the presentation generator
    (## Slide N: <title> followed by bullets) into a python-pptx Presentation.
    """
    import io
    import re
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor

    pres = Presentation()
    pres.slide_width = Inches(13.333)
    pres.slide_height = Inches(7.5)

    blank_layout = pres.slide_layouts[6]  # truly blank

    def add_slide(slide_title: str, bullets: List[str], notes: str = "") -> None:
        slide = pres.slides.add_slide(blank_layout)
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12.0), Inches(1.0))
        tf = title_box.text_frame
        tf.text = slide_title
        for p in tf.paragraphs:
            p.alignment = PP_ALIGN.LEFT
            for r in p.runs:
                r.font.size = Pt(28)
                r.font.bold = True
                r.font.color.rgb = RGBColor(0x0F, 0x17, 0x2A)
        # Bullets
        body_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.7), Inches(11.9), Inches(5.0))
        bf = body_box.text_frame
        bf.word_wrap = True
        if not bullets:
            bullets = [""]
        for i, b in enumerate(bullets):
            para = bf.paragraphs[0] if i == 0 else bf.add_paragraph()
            para.text = b
            para.level = 0
            for r in para.runs:
                r.font.size = Pt(18)
                r.font.color.rgb = RGBColor(0x37, 0x41, 0x51)
        # Speaker notes
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    lines = markdown.replace("\r\n", "\n").split("\n")
    current_title = title
    current_bullets: List[str] = []
    current_notes = ""
    slide_count = 0
    # Accept any separator after the slide number (": ", " - ", en/em dash) or none.
    slide_heading_re = re.compile(r"^##\s*Slide\s*\d+\s*[:\-–—]?\s*(.*)$", re.IGNORECASE)
    bullet_re = re.compile(r"^[-*]\s+(.+)$")
    notes_re = re.compile(r"^\s*Notes:\s*(.+)$", re.IGNORECASE)

    def flush() -> None:
        nonlocal current_bullets, current_notes, slide_count
        if current_bullets or current_notes or slide_count == 0:
            # Strip inline markdown emphasis chars for the PPTX text
            clean = [re.sub(r"\*\*|__|`", "", b).strip() for b in current_bullets if b.strip()]
            add_slide(current_title, clean, current_notes)
            slide_count += 1
        current_bullets = []
        current_notes = ""

    for raw in lines:
        line = raw.rstrip()
        if not line.strip():
            continue
        m_h = slide_heading_re.match(line)
        if m_h:
            flush()
            current_title = m_h.group(1).strip().rstrip(".")
            continue
        m_n = notes_re.match(line)
        if m_n:
            current_notes = (current_notes + " " + m_n.group(1).strip()).strip()
            continue
        m_b = bullet_re.match(line)
        if m_b:
            current_bullets.append(m_b.group(1).strip())
            continue
        # plain line: treat as a bullet so nothing is dropped
        if line.startswith("#"):
            # ignore stray top-level headings
            continue
        current_bullets.append(line.strip())

    # Final slide
    if current_bullets or current_notes:
        flush()
    # Safety: at least one slide
    if slide_count == 0:
        add_slide(title, ["(empty deck)"], "")

    buffer = io.BytesIO()
    pres.save(buffer)
    return buffer.getvalue()


def _render_pptx_poster(title: str, markdown: str) -> bytes:
    """
    Render an event poster as a single large-format (A3 portrait) PPTX slide:
    an optional logo (from an embedded <img src="/clients/...">), a big title,
    a subtitle, and the body sections. Posters are visual single-slide documents,
    so they use their own renderer rather than the slide-deck parser.
    """
    import io
    import re
    from pathlib import Path
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor

    pres = Presentation()
    pres.slide_width = Inches(11.69)   # A3 portrait
    pres.slide_height = Inches(16.54)
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    sw = pres.slide_width

    content = markdown.replace("\r\n", "\n")
    green = RGBColor(0x2E, 0x7D, 0x32)
    dark = RGBColor(0x0F, 0x17, 0x2A)
    grey = RGBColor(0x37, 0x41, 0x51)

    # Optional logo from an embedded <img src="/clients/xxx.png">
    top = Inches(0.7)
    m_img = re.search(r'<img[^>]+src="([^"]+)"', content)
    if m_img and m_img.group(1).startswith("/clients/"):
        repo_root = Path(__file__).resolve().parents[2]
        logo_path = repo_root / "frontend" / "public" / m_img.group(1).lstrip("/")
        if logo_path.exists():
            try:
                pic = slide.shapes.add_picture(str(logo_path), Inches(0), top, height=Inches(1.2))
                pic.left = int((sw - pic.width) / 2)
                top = Inches(2.3)
            except Exception:
                pass

    # Strip HTML, then classify lines.
    text = re.sub(r"<[^>]+>", "", content)
    big_title = None
    subtitle = ""
    body: List[tuple] = []  # (kind, text): 'h' heading, 'b' bullet, 'p' paragraph
    for raw in text.split("\n"):
        s = raw.strip()
        if not s or s == "---":
            continue
        if s.startswith("# ") and big_title is None:
            big_title = s[2:].strip()
            continue
        if s.startswith("## ") and not subtitle:
            subtitle = s[3:].strip()
            continue
        if s.startswith("### "):
            body.append(("h", s[4:].strip()))
            continue
        if s.startswith(("- ", "* ")):
            body.append(("b", s[2:].strip()))
            continue
        if s.startswith("#"):
            body.append(("p", s.lstrip("#").strip()))
            continue
        body.append(("p", s))

    def clean(x: str) -> str:
        return re.sub(r"\*\*|__|`|\*", "", x).strip()

    # Title
    tb = slide.shapes.add_textbox(Inches(0.5), top, sw - Inches(1.0), Inches(1.6))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.text = clean(big_title or title)
    for p in tf.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        for r in p.runs:
            r.font.size = Pt(40)
            r.font.bold = True
            r.font.color.rgb = green
    top = top + Inches(1.7)

    # Subtitle
    if subtitle:
        sb = slide.shapes.add_textbox(Inches(0.8), top, sw - Inches(1.6), Inches(1.0))
        sf = sb.text_frame
        sf.word_wrap = True
        sf.text = clean(subtitle)
        for p in sf.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                r.font.size = Pt(22)
                r.font.color.rgb = dark
        top = top + Inches(1.1)

    # Body
    bb = slide.shapes.add_textbox(Inches(1.0), top, sw - Inches(2.0), pres.slide_height - top - Inches(0.6))
    bf = bb.text_frame
    bf.word_wrap = True
    first = True
    for kind, txt in body:
        para = bf.paragraphs[0] if first else bf.add_paragraph()
        first = False
        if kind == "b":
            para.text = "• " + clean(txt)
            sz, bold, col = 16, False, grey
        elif kind == "h":
            para.text = clean(txt)
            para.space_before = Pt(10)
            sz, bold, col = 20, True, dark
        else:
            para.text = clean(txt)
            sz, bold, col = 16, False, grey
        for r in para.runs:
            r.font.size = Pt(sz)
            r.font.bold = bold
            r.font.color.rgb = col

    buffer = io.BytesIO()
    pres.save(buffer)
    return buffer.getvalue()


@router.get(
    "/{document_id}/export-pptx",
    summary="Export presentation or poster as PPTX",
    description="Render a saved presentation (slides) or event poster (single A3 slide) as a real .pptx file",
)
async def export_document_pptx(
    document_id: UUID,
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db),
):
    document = db.query(UserDocument).filter(
        and_(
            UserDocument.id == document_id,
            UserDocument.user_id == current_user.id,
        )
    ).first()
    if not document:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Document not found")

    try:
        tags = document.tags or []
        if "event_poster" in tags:
            payload = _render_pptx_poster(document.title or "Poster", document.content or "")
        else:
            payload = _render_pptx_from_slides_markdown(document.title or "Presentation", document.content or "")
    except Exception as exc:
        logger.exception("Failed to render PPTX for %s", document_id)
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=f"Failed to export PPTX: {exc}",
        ) from exc

    import io
    filename = _safe_filename(document.title or "presentation", "pptx")
    return StreamingResponse(
        io.BytesIO(payload),
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f'attachment; filename="{filename}"'},
    )
