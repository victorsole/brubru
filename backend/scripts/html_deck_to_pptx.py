"""
Convert a Brubru slide-deck HTML (fixed 1280x720 <section class="slide"> blocks)
into a PowerPoint: Playwright screenshots each slide at @2x, python-pptx drops
them full-bleed onto 13.333in x 7.5in slides.

Usage: python3.12 html_deck_to_pptx.py <deck.html> <out.pptx>
"""
import sys
from pathlib import Path

from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.util import Inches

def main(html_path: str, out_path: str) -> None:
    html = Path(html_path).resolve()
    out = Path(out_path).resolve()
    shots_dir = out.parent / (out.stem + "_slides")
    shots_dir.mkdir(exist_ok=True)

    pngs = []
    with sync_playwright() as p:
        browser = p.chromium.launch()
        page = browser.new_page(viewport={"width": 1280, "height": 720}, device_scale_factor=2)
        page.goto(html.as_uri())
        page.wait_for_load_state("networkidle")
        page.wait_for_timeout(800)  # let fonts/images settle
        slides = page.locator("section.slide")
        n = slides.count()
        for i in range(n):
            el = slides.nth(i)
            el.scroll_into_view_if_needed()
            page.wait_for_timeout(200)
            png = shots_dir / f"slide_{i+1:02d}.png"
            el.screenshot(path=str(png))
            pngs.append(png)
            print(f"[OK] captured slide {i+1}/{n} -> {png.name}")
        browser.close()

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for png in pngs:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(str(png), 0, 0, width=prs.slide_width, height=prs.slide_height)
    prs.save(str(out))
    print(f"[OK] wrote {out} ({out.stat().st_size/1024:.0f} KB, {len(pngs)} slides)")

if __name__ == "__main__":
    main(sys.argv[1], sys.argv[2])
